"""Utilities for loading and normalising documents."""
from __future__ import annotations

import io
from pathlib import Path
from typing import Iterable, Iterator, Tuple

import docx2txt
from pptx import Presentation
from PyPDF2 import PdfReader
try:  # pragma: no cover - optional dependencies
    import docx2txt
except Exception:  # noqa: BLE001 - optional dependency
    docx2txt = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependencies
    from pptx import Presentation
except Exception:  # noqa: BLE001 - optional dependency
    Presentation = None  # type: ignore[assignment]

try:  # pragma: no cover - optional dependencies
    from PyPDF2 import PdfReader
except Exception:  # noqa: BLE001 - optional dependency
    PdfReader = None  # type: ignore[assignment]

SUPPORTED_EXTENSIONS: Tuple[str, ...] = (".txt", ".md", ".pdf", ".docx", ".pptx")


class UnsupportedDocumentError(ValueError):
    """Raised when a document has an unsupported extension."""


def _load_pdf(path: Path) -> str:
    if PdfReader is None:
        raise UnsupportedDocumentError("PyPDF2 is required to load PDF files")
    reader = PdfReader(str(path))
    text_buffer = io.StringIO()
    for page in reader.pages:
        text_buffer.write(page.extract_text() or "")
        text_buffer.write("\n")
    return text_buffer.getvalue()


def _load_docx(path: Path) -> str:
    if docx2txt is None:
        raise UnsupportedDocumentError("docx2txt is required to load DOCX files")
    return docx2txt.process(str(path)) or ""


def _load_pptx(path: Path) -> str:
    if Presentation is None:
        raise UnsupportedDocumentError("python-pptx is required to load PPTX files")
    presentation = Presentation(str(path))
    text_runs: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)


def _load_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="ignore")


_LOADERS = {
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".txt": _load_text,
    ".md": _load_text,
}


def load_document(path: Path) -> str:
    """Load a document and return its textual representation."""
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        raise UnsupportedDocumentError(
            f"Unsupported document type '{suffix}'. Supported: {SUPPORTED_EXTENSIONS}"
        )
    loader = _LOADERS[suffix]
    return loader(path)


def iter_supported_files(paths: Iterable[Path]) -> Iterator[Path]:
    """Yield only supported files from an iterable of paths."""
    for path in paths:
        if path.suffix.lower() in SUPPORTED_EXTENSIONS:
            yield path
